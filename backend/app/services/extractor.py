from pathlib import Path

from pptx import Presentation


def extract_deck(file_path: Path, deck_id: str) -> list[dict]:
    """Open a .pptx and extract slide metadata (title, dimensions)."""
    prs = Presentation(str(file_path))

    # Slide dimensions in EMU -> pixels at 96 DPI
    width_px = int(prs.slide_width / 914400 * 96)
    height_px = int(prs.slide_height / 914400 * 96)

    slides_info: list[dict] = []

    for idx, slide in enumerate(prs.slides):
        title = None
        for shape in slide.shapes:
            if (
                shape.has_text_frame
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx in (0, 1)
            ):
                title = shape.text_frame.text.strip()
                if title:
                    break

        slides_info.append(
            {
                "slide_index": idx,
                "title": title or None,
                "width_px": width_px,
                "height_px": height_px,
            }
        )

    return slides_info
